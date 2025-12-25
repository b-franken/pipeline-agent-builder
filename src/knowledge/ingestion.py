import hashlib
import logging
import os
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, BinaryIO, Final

from langchain_core.documents import Document

from src.knowledge.chunking import ChunkConfig, ChunkingStrategy, chunk_document

if TYPE_CHECKING:
    from src.memory.vector_store import VectorMemory
    from src.storage.repository import StorageRepository

logger: Final = logging.getLogger(__name__)

MAX_FILE_SIZE_BYTES: Final = 50 * 1024 * 1024
MAX_CONTENT_LENGTH: Final = 10_000_000


@dataclass
class IngestResult:
    document_id: str
    filename: str
    chunks_created: int
    total_tokens_estimate: int
    metadata: dict[str, object] = field(default_factory=dict)
    success: bool = True
    error: str | None = None


@dataclass
class DocumentRecord:
    document_id: str
    filename: str
    content_hash: str
    chunk_count: int
    ingested_at: datetime
    metadata: dict[str, object]


class DocumentIngester:
    def __init__(
        self,
        vector_store: VectorMemory | None = None,
        chunk_config: ChunkConfig | None = None,
        collection_name: str = "knowledge_base",
        max_file_size: int = MAX_FILE_SIZE_BYTES,
        max_content_length: int = MAX_CONTENT_LENGTH,
    ) -> None:
        self._vector_store: VectorMemory | None = vector_store
        self._chunk_config = chunk_config or ChunkConfig()
        self._collection_name = collection_name
        self._max_file_size = max_file_size
        self._max_content_length = max_content_length
        self._documents: dict[str, DocumentRecord] = {}
        self._repo: StorageRepository | None = None

    async def _get_repo(self) -> StorageRepository:
        if self._repo is None:
            from src.storage import get_repository

            self._repo = await get_repository()
        return self._repo

    async def _check_duplicate_async(self, content_hash: str) -> DocumentRecord | None:
        repo = await self._get_repo()
        existing = await repo.find_document_by_hash(content_hash)
        if existing:
            return DocumentRecord(
                document_id=existing.id,
                filename=existing.filename,
                content_hash=existing.content_hash,
                chunk_count=existing.chunk_count,
                ingested_at=existing.created_at or datetime.now(),
                metadata={},
            )
        return None

    def _get_vector_store(self) -> VectorMemory:
        if self._vector_store is None:
            from src.memory.vector_store import VectorMemory

            self._vector_store = VectorMemory(collection_name=self._collection_name)
        return self._vector_store

    def _content_hash(self, content: str) -> str:
        return hashlib.sha256(content.encode()).hexdigest()

    def _generate_doc_id(self, content_hash: str) -> str:
        return content_hash[:16]

    def _validate_file_size(self, file_path: Path) -> None:
        size = file_path.stat().st_size
        if size > self._max_file_size:
            raise ValueError(f"File size {size} bytes exceeds maximum {self._max_file_size} bytes")

    def _validate_content_length(self, content: str) -> None:
        if len(content) > self._max_content_length:
            raise ValueError(f"Content length {len(content)} exceeds maximum {self._max_content_length}")

    def _parse_file(self, file_path: Path) -> tuple[str, dict[str, object]]:
        self._validate_file_size(file_path)
        suffix = file_path.suffix.lower()
        metadata: dict[str, object] = {
            "filename": file_path.name,
            "file_type": suffix,
            "file_size": file_path.stat().st_size,
        }

        if suffix == ".pdf":
            return self._parse_pdf(file_path), metadata
        elif suffix == ".docx":
            return self._parse_docx(file_path), metadata
        elif suffix == ".xlsx":
            return self._parse_xlsx(file_path), metadata
        elif suffix == ".pptx":
            return self._parse_pptx(file_path), metadata
        elif suffix == ".csv":
            return self._parse_csv(file_path), metadata
        elif suffix in (".md", ".markdown") or suffix in (".txt", ".text"):
            return file_path.read_text(encoding="utf-8"), metadata
        elif suffix in (".py", ".js", ".ts", ".java", ".go", ".rs", ".cpp", ".c", ".rb", ".jsx", ".tsx"):
            content = file_path.read_text(encoding="utf-8")
            metadata["language"] = suffix[1:]  # Remove the dot
            return content, metadata
        elif suffix == ".json":
            return file_path.read_text(encoding="utf-8"), metadata
        elif suffix == ".html":
            return self._parse_html(file_path), metadata
        else:
            try:
                return file_path.read_text(encoding="utf-8"), metadata
            except UnicodeDecodeError as e:
                raise ValueError(f"Cannot parse binary file: {file_path.name}") from e

    def _parse_pdf(self, file_path: Path) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return "\n\n".join(text_parts)
        except ImportError:
            pass

        try:
            import pdfplumber

            with pdfplumber.open(str(file_path)) as pdf:
                text_parts = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return "\n\n".join(text_parts)
        except ImportError as e:
            raise ImportError("PDF parsing requires pypdf or pdfplumber. Install with: pip install pypdf") from e

    def _parse_docx(self, file_path: Path) -> str:
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(str(file_path))
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs)
        except ImportError as e:
            raise ImportError(
                "Word document parsing requires python-docx. Install with: pip install python-docx"
            ) from e

    def _parse_xlsx(self, file_path: Path) -> str:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                all_text.append(f"## Sheet: {sheet_name}\n")
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        rows.append(" | ".join(row_values))
                all_text.append("\n".join(rows))
            wb.close()
            return "\n\n".join(all_text)
        except ImportError as e:
            raise ImportError("Excel parsing requires openpyxl. Install with: pip install openpyxl") from e

    def _parse_pptx(self, file_path: Path) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"## Slide {i}"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                if len(slide_content) > 1:
                    slides_text.append("\n".join(slide_content))
            return "\n\n".join(slides_text)
        except ImportError as e:
            raise ImportError("PowerPoint parsing requires python-pptx. Install with: pip install python-pptx") from e

    def _parse_csv(self, file_path: Path) -> str:
        import csv

        with open(file_path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(row))
            return "\n".join(rows)

    def _parse_html(self, file_path: Path) -> str:
        try:
            from bs4 import BeautifulSoup

            html_content = file_path.read_text(encoding="utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator="\n")
            lines = (line.strip() for line in text.splitlines())
            return "\n".join(line for line in lines if line)
        except ImportError:
            import re

            html_content = file_path.read_text(encoding="utf-8")
            flags = re.DOTALL | re.IGNORECASE
            html_content = re.sub(r"<script\b[^>]*>.*?</script\b[^>]*>", "", html_content, flags=flags)
            html_content = re.sub(r"<style\b[^>]*>.*?</style\b[^>]*>", "", html_content, flags=flags)
            text = re.sub(r"<[^>]+>", " ", html_content)
            text = re.sub(r"\s+", " ", text)
            return text.strip()

    def _store_chunks(self, chunks: list[Document], document_id: str) -> int:
        vs = self._get_vector_store()
        for chunk in chunks:
            chunk.metadata["document_id"] = document_id
        vectorstore = vs._get_vectorstore()
        if vectorstore is not None:
            vectorstore.add_documents(chunks)
        return len(chunks)

    def ingest_file(
        self,
        file_path: str | Path,
        metadata: dict[str, object] | None = None,
        strategy: ChunkingStrategy | None = None,
    ) -> IngestResult:
        path = Path(file_path)
        if not path.exists():
            return IngestResult(
                document_id="",
                filename=str(file_path),
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error=f"File not found: {file_path}",
            )

        try:
            content, file_metadata = self._parse_file(path)
            combined_metadata = {**file_metadata, **(metadata or {})}

            lang = file_metadata.get("language")
            return self.ingest_text(
                content=content,
                filename=path.name,
                metadata=combined_metadata,
                strategy=strategy,
                language=str(lang) if lang else None,
            )
        except Exception as e:
            return IngestResult(
                document_id="",
                filename=path.name,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error=str(e),
            )

    async def ingest_text_async(
        self,
        content: str,
        filename: str = "document.txt",
        metadata: dict[str, object] | None = None,
        strategy: ChunkingStrategy | None = None,
        language: str | None = None,
    ) -> IngestResult:
        if not content.strip():
            return IngestResult(
                document_id="",
                filename=filename,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error="Empty content",
            )

        try:
            self._validate_content_length(content)
            content_hash = self._content_hash(content)
            document_id = self._generate_doc_id(content_hash)

            existing = await self._check_duplicate_async(content_hash)
            if existing:
                return IngestResult(
                    document_id=existing.document_id,
                    filename=filename,
                    chunks_created=existing.chunk_count,
                    total_tokens_estimate=0,
                    metadata={"duplicate_of": existing.document_id, "content_hash": content_hash},
                    success=True,
                )

            config = ChunkConfig(
                chunk_size=self._chunk_config.chunk_size,
                chunk_overlap=self._chunk_config.chunk_overlap,
                strategy=strategy or self._chunk_config.strategy,
            )

            chunks = chunk_document(
                content=content,
                metadata=metadata or {},
                config=config,
                filename=filename,
                language=language,
            )

            stored_count = self._store_chunks(chunks, document_id)

            self._documents[document_id] = DocumentRecord(
                document_id=document_id,
                filename=filename,
                content_hash=content_hash,
                chunk_count=stored_count,
                ingested_at=datetime.now(),
                metadata=metadata or {},
            )

            total_chars = sum(len(c.page_content) for c in chunks)
            token_estimate = total_chars // 4

            return IngestResult(
                document_id=document_id,
                filename=filename,
                chunks_created=stored_count,
                total_tokens_estimate=token_estimate,
                metadata={"content_hash": content_hash},
                success=True,
            )

        except Exception as e:
            return IngestResult(
                document_id="",
                filename=filename,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error=str(e),
            )

    def ingest_text(
        self,
        content: str,
        filename: str = "document.txt",
        metadata: dict[str, object] | None = None,
        strategy: ChunkingStrategy | None = None,
        language: str | None = None,
    ) -> IngestResult:
        if not content.strip():
            return IngestResult(
                document_id="",
                filename=filename,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error="Empty content",
            )

        try:
            self._validate_content_length(content)
            content_hash = self._content_hash(content)
            document_id = self._generate_doc_id(content_hash)

            for doc in self._documents.values():
                if doc.content_hash == content_hash:
                    return IngestResult(
                        document_id=doc.document_id,
                        filename=filename,
                        chunks_created=doc.chunk_count,
                        total_tokens_estimate=0,
                        metadata={"duplicate_of": doc.document_id, "content_hash": content_hash},
                        success=True,
                    )

            config = ChunkConfig(
                chunk_size=self._chunk_config.chunk_size,
                chunk_overlap=self._chunk_config.chunk_overlap,
                strategy=strategy or self._chunk_config.strategy,
            )

            chunks = chunk_document(
                content=content,
                metadata=metadata or {},
                config=config,
                filename=filename,
                language=language,
            )

            stored_count = self._store_chunks(chunks, document_id)

            self._documents[document_id] = DocumentRecord(
                document_id=document_id,
                filename=filename,
                content_hash=content_hash,
                chunk_count=stored_count,
                ingested_at=datetime.now(),
                metadata=metadata or {},
            )

            total_chars = sum(len(c.page_content) for c in chunks)
            token_estimate = total_chars // 4

            return IngestResult(
                document_id=document_id,
                filename=filename,
                chunks_created=stored_count,
                total_tokens_estimate=token_estimate,
                metadata={"content_hash": content_hash},
                success=True,
            )

        except Exception as e:
            return IngestResult(
                document_id="",
                filename=filename,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error=str(e),
            )

    def ingest_upload(
        self,
        file_data: bytes | BinaryIO,
        filename: str,
        metadata: dict[str, object] | None = None,
    ) -> IngestResult:
        data = bytes(file_data.read()) if hasattr(file_data, "read") else bytes(file_data)

        if len(data) > self._max_file_size:
            return IngestResult(
                document_id="",
                filename=filename,
                chunks_created=0,
                total_tokens_estimate=0,
                success=False,
                error=f"File size {len(data)} bytes exceeds maximum {self._max_file_size} bytes",
            )

        suffix = Path(filename).suffix.lower()
        safe_suffixes = {".pdf": ".pdf", ".docx": ".docx", ".xlsx": ".xlsx", ".pptx": ".pptx"}

        if suffix in safe_suffixes:
            import tempfile

            safe_suffix = safe_suffixes[suffix]
            with tempfile.NamedTemporaryFile(suffix=safe_suffix, delete=False) as tmp:
                tmp.write(data)
                tmp_path = tmp.name

            try:
                result = self.ingest_file(tmp_path, metadata)
                result.filename = filename  # Use original filename
                return result
            finally:
                os.unlink(tmp_path)
        else:
            try:
                content = data.decode("utf-8")
            except UnicodeDecodeError:
                return IngestResult(
                    document_id="",
                    filename=filename,
                    chunks_created=0,
                    total_tokens_estimate=0,
                    success=False,
                    error="Cannot decode file as text. Supported binary formats: PDF, DOCX, XLSX, PPTX",
                )

            return self.ingest_text(content, filename, metadata)

    async def sync_from_database(self) -> int:
        repo = await self._get_repo()
        docs = await repo.list_documents(active_only=True)
        synced = 0
        for doc in docs:
            if doc.id not in self._documents:
                self._documents[doc.id] = DocumentRecord(
                    document_id=doc.id,
                    filename=doc.filename,
                    content_hash=doc.content_hash,
                    chunk_count=doc.chunk_count,
                    ingested_at=doc.created_at or datetime.now(),
                    metadata={},
                )
                synced += 1
        return synced

    def list_documents(self) -> list[DocumentRecord]:
        return list(self._documents.values())

    def get_document(self, document_id: str) -> DocumentRecord | None:
        return self._documents.get(document_id)

    def delete_document(self, document_id: str) -> bool:
        if document_id not in self._documents:
            return False

        try:
            vs = self._get_vector_store()
            vectorstore = vs._get_vectorstore()
            if vectorstore is not None:
                collection = vectorstore._collection
                collection.delete(where={"document_id": document_id})
            del self._documents[document_id]
            return True
        except Exception:
            logger.exception("Failed to delete document %s", document_id)
            return False

    def clear_all(self) -> int:
        count = len(self._documents)

        try:
            vs = self._get_vector_store()
            vectorstore = vs._get_vectorstore()
            if vectorstore is not None:
                collection = vectorstore._collection
                if collection.count() > 0:
                    collection.delete(where={})
            self._documents.clear()
            return count
        except Exception:
            logger.exception("Failed to clear all documents")
            self._documents.clear()
            return count
